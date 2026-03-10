"""
extract_text.py — Extracts text content from PDF and PPTX files.

Handles three scenarios:
  1. Normal selectable text  → extracted directly (fast)
  2. Text embedded as images → OCR via EasyOCR (GPU-accelerated)
  3. Non-text visuals        → detected and described in context

Requires:  pip install PyMuPDF python-pptx easyocr Pillow
"""
import io
import os
import sys
import logging

log = logging.getLogger("extract")

# Minimum characters per page to consider text extraction "successful"
MIN_CHARS_PER_PAGE = 20


# ──────────────────────────────────────────────
# OCR engine (lazy-loaded, singleton)
# ──────────────────────────────────────────────

_ocr_reader = None


def _get_ocr_reader():
    """Lazy-load EasyOCR reader (uses GPU if available)."""
    global _ocr_reader
    if _ocr_reader is None:
        import easyocr
        _ocr_reader = easyocr.Reader(
            ["en"],
            gpu=True,       # will use your A100 automatically
            verbose=False,
        )
        log.info("EasyOCR reader initialized (GPU-accelerated)")
    return _ocr_reader


def ocr_image(image_bytes: bytes) -> str:
    """Run OCR on an image (as bytes) and return the extracted text."""
    from PIL import Image
    import numpy as np

    reader = _get_ocr_reader()
    img = Image.open(io.BytesIO(image_bytes)).convert("RGB")
    img_array = np.array(img)

    results = reader.readtext(img_array, detail=0, paragraph=True)
    return "\n".join(results)


# ──────────────────────────────────────────────
# Image / diagram detection
# ──────────────────────────────────────────────

def _describe_and_save_page_images(page, page_num: int, source_path: str) -> list[str]:
    """Detect embedded images on a PDF page, save them if they are useful, and return descriptions."""
    descriptions = []
    image_list = page.get_images(full=True)

    if not image_list:
        return descriptions

    # Ensure assets directory exists
    base_dir = os.path.dirname(source_path)
    base_name = os.path.splitext(os.path.basename(source_path))[0]
    assets_dir = os.path.join(base_dir, "assets")

    for img_index, img_info in enumerate(image_list, start=1):
        xref = img_info[0]
        try:
            base_image = page.parent.extract_image(xref)
            width = base_image.get("width", 0)
            height = base_image.get("height", 0)

            # ── Filtering Heuristics ──
            # Skip tiny images (icons, bullets, decorations)
            if width < 150 or height < 150:
                continue
            
            # Skip extreme aspect ratios (thin lines, borders)
            aspect = width / max(height, 1)
            if aspect < 0.2 or aspect > 5.0:
                continue

            # Classify by aspect ratio and size
            area = width * height
            if area > 200_000:
                kind = "large diagram/figure"
            elif 0.8 < aspect < 1.2:
                kind = "square image (possibly a photo or chart)"
            elif aspect > 2.0:
                kind = "wide image (possibly a banner, table, or flowchart)"
            elif aspect < 0.5:
                kind = "tall image (possibly a column chart or listing)"
            else:
                kind = "image"
                
            # ── Save the image ──
            os.makedirs(assets_dir, exist_ok=True)
            ext = base_image.get("ext", "png")
            img_filename = f"{base_name}_page{page_num}_img{img_index}.{ext}"
            img_path = os.path.join(assets_dir, img_filename)
            
            with open(img_path, "wb") as f:
                f.write(base_image["image"])

            descriptions.append(
                f"[Visual Feature: {kind}, {width}x{height}px — saved at assets/{img_filename}]"
            )
        except Exception as e:
            descriptions.append(f"[Visual: embedded image (could not inspect: {e})]")

    return descriptions


# ──────────────────────────────────────────────
# PDF extraction (selectable text + OCR fallback)
# ──────────────────────────────────────────────

def extract_from_pdf(filepath: str) -> str:
    """Extract text from a PDF. Falls back to OCR for image-based pages."""
    import fitz  # PyMuPDF

    text_parts = []
    ocr_pages = 0

    with fitz.open(filepath) as doc:
        total_pages = len(doc)
        log.info(f"  PDF has {total_pages} pages")

        for page_num, page in enumerate(doc, start=1):
            page_header = f"--- Page {page_num} ---"

            # ── Step 1: Try normal text extraction ──
            page_text = page.get_text("text").strip()

            # ── Step 2: Detect, save, and annotate visual elements ──
            image_descriptions = _describe_and_save_page_images(page, page_num, filepath)

            # ── Step 3: If text is too sparse, try OCR ──
            if len(page_text) < MIN_CHARS_PER_PAGE:
                log.info(f"  Page {page_num}: only {len(page_text)} chars → running OCR")
                try:
                    # Render the full page as an image at 300 DPI for OCR
                    pix = page.get_pixmap(dpi=300)
                    img_bytes = pix.tobytes("png")
                    ocr_text = ocr_image(img_bytes)

                    if ocr_text.strip():
                        page_text = ocr_text.strip()
                        ocr_pages += 1
                    else:
                        page_text = "[No readable text detected on this page]"
                except Exception as e:
                    log.warning(f"  Page {page_num}: OCR failed — {e}")
                    if not page_text:
                        page_text = "[OCR failed for this page]"

            # ── Combine text + visual descriptions ──
            parts = [page_header, page_text]
            if image_descriptions:
                parts.append("\n".join(image_descriptions))

            text_parts.append("\n".join(parts))

    if ocr_pages:
        log.info(f"  OCR was used on {ocr_pages}/{total_pages} pages")

    return "\n\n".join(text_parts)


# ──────────────────────────────────────────────
# PPTX extraction
# ──────────────────────────────────────────────

def extract_from_pptx(filepath: str) -> str:
    """Extract text from all slides (shapes + notes) of a PPTX file."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    text_parts = []
    prs = Presentation(filepath)
    
    base_dir = os.path.dirname(filepath)
    base_name = os.path.splitext(os.path.basename(filepath))[0]
    assets_dir = os.path.join(base_dir, "assets")

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_texts = []
        image_count = 0

        for shape_idx, shape in enumerate(slide.shapes, start=1):
            # Text content
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    line = paragraph.text.strip()
                    if line:
                        slide_texts.append(line)

            # Tables
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_text = " | ".join(
                        cell.text.strip() for cell in row.cells if cell.text.strip()
                    )
                    if row_text:
                        slide_texts.append(f"[Table row]: {row_text}")

            # Images / media
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    width = image.size[0]
                    height = image.size[1]
                    
                    # Basic heuristics to skip tiny icons
                    if width > 1000000 and height > 1000000: # ~150px in EMU
                        os.makedirs(assets_dir, exist_ok=True)
                        ext = image.ext
                        img_filename = f"{base_name}_slide{slide_num}_img{shape_idx}.{ext}"
                        img_path = os.path.join(assets_dir, img_filename)
                        
                        with open(img_path, "wb") as f:
                            f.write(image.blob)
                        
                        slide_texts.append(f"[Visual Feature: Diagram/Picture — saved at assets/{img_filename}]")
                        image_count += 1
                except Exception as e:
                    slide_texts.append(f"[Visual: embedded image (could not extract: {e})]")
                    
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                # Groups often contain diagrams
                slide_texts.append("[Visual: grouped diagram/figure]")

        # Speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_texts.append(f"[Speaker Notes]: {notes}")

        if slide_texts:
            text_parts.append(
                f"--- Slide {slide_num} ---\n" + "\n".join(slide_texts)
            )

    return "\n\n".join(text_parts)


# ──────────────────────────────────────────────
# Dispatcher
# ──────────────────────────────────────────────

def extract_text(filepath: str) -> str:
    """Dispatch to the correct extractor based on file extension."""
    ext = os.path.splitext(filepath)[1].lower()
    if ext == ".pdf":
        return extract_from_pdf(filepath)
    elif ext in (".pptx", ".ppt"):
        return extract_from_pptx(filepath)
    else:
        raise ValueError(f"Unsupported file type: {ext}")


# ──────────────────────────────────────────────
# CLI
# ──────────────────────────────────────────────

if __name__ == "__main__":
    logging.basicConfig(
        level=logging.INFO,
        format="%(asctime)s  %(levelname)-8s  %(message)s",
        datefmt="%H:%M:%S",
    )

    if len(sys.argv) < 2:
        print("Usage: python extract_text.py <filepath>")
        sys.exit(1)

    text = extract_text(sys.argv[1])
    print(text[:5000])
    print(f"\n... (total {len(text)} characters extracted)")
